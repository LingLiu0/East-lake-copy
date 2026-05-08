#!/usr/bin/env python3
"""
Obsidian AI 工作流 - Karpathy llm-wiki 方法实现

核心操作：
- Ingest: 摄入资料 → 自动生成摘要、提取概念、建立链接
- Query: 问答 → 答案沉淀到 wiki
- Lint: 健康检查 → 发现孤立页面、矛盾、过时内容
"""
import argparse
import json
import os
import re
import subprocess
import sys
import hashlib
from datetime import datetime
from pathlib import Path
from typing import Optional, Set

VERSION = "3.0"

# ============================================================
# 核心配置 (Karpathy llm-wiki 方法)
# ============================================================

VAULT = Path(os.getcwd())

# llm-wiki 目录结构
RAW = VAULT / "raw"           # 原始资料（不可变）
WIKI = VAULT / "wiki"         # LLM 编译产物
OUTPUTS = VAULT / "outputs"   # 运行时输出
ARCHIVE = VAULT / "_archive"  # 归档
CLIPPINGS = RAW / "clippings"  # 网页剪藏
ARTICLES = RAW / "articles"    # 文章/文档

# 支持的文件格式
SUPPORTED_EXTENSIONS = ['.md', '.txt', '.pdf', '.docx', '.ppt', '.pptx', '.html', '.htm']

# 兼容旧目录
INBOX = VAULT / "_inbox" if (VAULT / "_inbox").exists() else RAW
ACHIEVEMENTS = WIKI

# 创建必要目录
for d in [RAW, CLIPPINGS, ARTICLES, WIKI, OUTPUTS, WIKI / "indexes", WIKI / "concepts", WIKI / "summaries", OUTPUTS / "qa", OUTPUTS / "health"]:
    d.mkdir(parents=True, exist_ok=True)


# ============================================================
# 工具函数
# ============================================================

def parse_file(file_path: Path) -> str:
    """解析不同格式的文件，返回文本内容"""
    ext = file_path.suffix.lower()
    content = ""

    try:
        if ext == '.md' or ext == '.txt':
            # Markdown 和纯文本直接读取
            content = file_path.read_text(encoding='utf-8', errors='ignore')

        elif ext == '.pdf':
            # PDF 解析
            try:
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        content += page.extract_text() + "\n"
            except ImportError:
                content = f"[PDF 文件 - 需要安装 PyPDF2: pip install PyPDF2]"

        elif ext == '.docx':
            # Word 文档解析
            try:
                import docx
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    content += para.text + "\n"
            except ImportError:
                content = f"[DOCX 文件 - 需要安装 python-docx: pip install python-docx]"

        elif ext in ['.ppt', '.pptx']:
            # PPT 解析
            try:
                import pptx
                prs = pptx.Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
            except ImportError:
                content = f"[PPT 文件 - 需要安装 python-pptx: pip install python-pptx]"

        elif ext in ['.html', '.htm']:
            # HTML 解析
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(file_path.read_text(encoding='utf-8', errors='ignore'), 'html.parser')
                for tag in soup(['script', 'style', 'nav', 'footer']):
                    tag.decompose()
                content = soup.get_text(separator='\n', strip=True)
            except ImportError:
                content = file_path.read_text(encoding='utf-8', errors='ignore')

        else:
            content = f"[不支持的格式: {ext}]"

    except Exception as e:
        content = f"[解析失败: {str(e)}]"

    return content


def compute_hash(content: bytes) -> str:
    return hashlib.md5(content).hexdigest()[:12]


def run_py(script: str, args: list = None):
    """运行 Python 脚本"""
    cmd = [sys.executable, str(VAULT / "scripts" / script)]
    if args:
        cmd.extend(args)
    return subprocess.run(cmd, cwd=str(VAULT)).returncode


def print_header(title: str):
    print(f"\n{'='*50}")
    print(f"  {title}")
    print(f"{'='*50}\n")


def print_status(items: list, empty_msg: str = "无"):
    if not items:
        print(f"  {empty_msg}")
        return
    for i, item in enumerate(items, 1):
        print(f"  {i}. {item}")


def extract_concepts(content: str) -> list[str]:
    """从内容中提取关键概念（简单版）"""
    # 常见概念模式
    patterns = [
        r'#+\s+(.{2,20})',           # 标题
        r'\*\*(.{2,20})\*\*',        # 加粗
        r'「([^」]{2,20})」',         # 中文引号
        r'"([^"]{2,20})"',           # 英文引号
    ]

    concepts = set()
    for pattern in patterns:
        matches = re.findall(pattern, content)
        concepts.update(m for m in matches if len(m) >= 2)

    return list(concepts)[:10]


def find_related_concepts(new_concept: str, existing_concepts: list[str]) -> list[str]:
    """找到相关概念"""
    related = []
    new_lower = new_concept.lower()
    for ec in existing_concepts:
        ec_lower = ec.lower()
        # 简单匹配：包含关系或共同词
        if new_lower in ec_lower or ec_lower in new_lower:
            related.append(ec)
        # 共同词
        new_words = set(new_lower.split())
        old_words = set(ec_lower.split())
        if new_words & old_words:
            related.append(ec)
    return related[:5]


# ============================================================
# 成果管理 (保留兼容)
# ============================================================

def cmd_achieve(args):
    """成果管理命令"""
    if args.action == "in" or args.action == "inbox":
        print_header("📥 原始资料箱")
        files = [f for f in RAW.rglob("*") if f.is_file() and not f.name.startswith('.')]
        print(f"\n  原始文件: {len(files)}")
        for f in files[:5]:
            print(f"    - {f.name}")
        print("\n  处理命令: python3 scripts/obsidian.py ai compile")

    elif args.action == "process" or args.action == "run":
        cmd_compile()

    elif args.action == "status":
        print_header("📊 系统状态")
        raw_count = len(list(RAW.rglob("*.md"))) if RAW.exists() else 0
        wiki_count = len(list(WIKI.rglob("*.md"))) if WIKI.exists() else 0
        print(f"  原始资料: {raw_count} 个")
        print(f"  Wiki: {wiki_count} 篇")


# ============================================================
# 网页收集
# ============================================================

def cmd_web(args):
    """网页收集命令"""
    if args.action == "add":
        if not args.url:
            print("错误: 请提供 URL")
            return
        url = args.url
        title = args.title or url.split('/')[-1][:50]

        web_inbox = RAW
        web_inbox.mkdir(parents=True, exist_ok=True)

        task = {
            "url": url,
            "title": title,
            "tags": args.tags.split(',') if args.tags else [],
            "added_at": datetime.now().isoformat(),
        }

        url_hash = compute_hash(url.encode())
        (web_inbox / f"{url_hash}.json").write_text(json.dumps(task, indent=2))
        print(f"  ✅ 已添加: {title}")
        print("  运行: python3 scripts/obsidian.py ai compile")

    elif args.action == "process":
        print_header("🌐 处理网页")
        processed = 0
        for task_file in RAW.glob("*.json"):
            task = json.loads(task_file.read_text())
            result = collect_web(task["url"], task.get("title"))
            if result["success"]:
                processed += 1
                task_file.unlink()
        print(f"\n  完成: {processed} 个网页")

    else:
        print_header("🌐 网页收集箱")
        files = list(RAW.glob("*.json"))
        print(f"  待处理: {len(files)} 个")


def collect_web(url: str, title: str = None) -> dict:
    """收集网页"""
    try:
        import requests
        from bs4 import BeautifulSoup

        resp = requests.get(url, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
        soup = BeautifulSoup(resp.text, 'html.parser')

        if not title:
            title = soup.title.string.strip()[:60] if soup.title else url

        for tag in soup(['script', 'style', 'nav', 'footer']):
            tag.decompose()

        content = soup.get_text(separator='\n', strip=True)
        lines = [l for l in content.split('\n') if len(l) > 20]
        content = '\n\n'.join(lines[:15])[:2000]

        # 保存到 raw/articles/
        safe_title = re.sub(r'[^\w\s-]', '', title).replace(' ', '-')[:40]
        date = datetime.now().strftime('%Y-%m-%d')
        md_file = RAW / "articles" / f"{date}-{safe_title}.md"
        md_file.parent.mkdir(parents=True, exist_ok=True)

        md_file.write_text(f"# {title}\n\n{content}", encoding='utf-8')
        return {"success": True, "title": title}
    except Exception as e:
        return {"success": False, "error": str(e)}


# ============================================================
# Karpathy llm-wiki 核心操作
# ============================================================

def cmd_ai(args):
    """AI 命令 - llm-wiki 方法"""
    action = args.action

    if action == "ask":
        if not args.query:
            print("错误: 请提供问题")
            return
        cmd_query(args.query)

    elif action == "chat":
        run_py("obsidian_chat.py", ["--interactive"])

    elif action == "stats":
        cmd_stats()

    elif action == "compile":
        cmd_compile(extract_concepts_flag=True)

    elif action == "lint":
        cmd_lint_enhanced()

    elif action == "index":
        cmd_update_index()

    elif action == "evolve":
        run_py("auto_evolve.py", ["--once"])

    elif action == "graph":
        """生成知识图谱"""
        run_py("generate_graph.py")


def cmd_stats():
    """显示知识库统计"""
    print_header("📊 知识库统计")

    raw_files = list(RAW.rglob("*.md")) if RAW.exists() else []
    wiki_files = list(WIKI.rglob("*.md")) if WIKI.exists() else []
    concepts = list((WIKI / "concepts").rglob("*.md")) if (WIKI / "concepts").exists() else []
    summaries = list((WIKI / "summaries").rglob("*.md")) if (WIKI / "summaries").exists() else []
    qa_outputs = list((OUTPUTS / "qa").rglob("*.md")) if (OUTPUTS / "qa").exists() else []

    print(f"  📂 原始资料: {len(raw_files)} 篇")
    print(f"  📚 Wiki: {len(wiki_files)} 篇")
    print(f"     - 概念: {len(concepts)} 个")
    print(f"     - 摘要: {len(summaries)} 篇")
    print(f"     - 问答沉淀: {len(qa_outputs)} 篇")

    # 链接统计
    total_links = 0
    for f in wiki_files:
        content = f.read_text(errors='ignore')
        total_links += len(re.findall(r'\[\[', content))
    print(f"     - 交叉链接: {total_links} 个")


def cmd_query(query: str, save_to_wiki: bool = True):
    """问答处理 - 包含结果沉淀"""
    print_header(f"🔍 查询: {query}")

    # 1. 搜索相关文档
    results = search_wiki(query)

    if not results:
        print("  未找到相关内容")
        print("  提示：放入更多资料到 raw/ 并运行 compile")
        return

    print(f"  找到 {len(results)} 个相关文档:\n")
    for i, r in enumerate(results[:5], 1):
        print(f"  {i}. {r['title']}")
        print(f"     {r['preview'][:80]}...")
        print()

    # 2. 如果启用 AI 且有 API Key，生成答案
    answer = None

    # 检查是否有自定义 API 或 Anthropic API
    use_custom_api = os.getenv("API_KEY") or os.getenv("ANTHROPIC_API_KEY")

    if use_custom_api:
        try:
            # 优先使用自定义 API
            if os.getenv("API_KEY"):
                from api_client import chat as custom_chat
                context = "\n\n".join([f"# {r['title']}\n{r['content'][:1000]}" for r in results[:3]])
                prompt = f"""基于以下知识库内容回答问题：{query}

{context}

请给出简洁准确的回答，并标注来源。"""
                answer = custom_chat(prompt)
            # 备用 Anthropic API
            elif os.getenv("ANTHROPIC_API_KEY"):
                import anthropic
                client = anthropic.Anthropic()
                context = "\n\n".join([f"# {r['title']}\n{r['content'][:1000]}" for r in results[:3]])
                prompt = f"""基于以下知识库内容回答问题：{query}

{context}

请给出简洁准确的回答，并标注来源。"""

                resp = client.messages.create(
                    model="claude-sonnet-4-20250514",
                    max_tokens=1000,
                    messages=[{"role": "user", "content": prompt}]
                )
                answer = resp.content[0].text

            if answer:
                print(f"  💡 AI 回答:\n{answer[:500]}...")
        except Exception as e:
            print(f"  ⚠️  AI 调用失败: {e}")

    # 3. 沉淀到 wiki（可选）
    if save_to_wiki and (results or answer):
        qa_file = OUTPUTS / "qa" / f"{datetime.now().strftime('%Y%m%d-%H%M%S')}-{query[:20]}.md"

        # 准备内容片段
        links = '\n'.join(f'- [[{r["path"]}]]' for r in results)
        ai_section = f'## AI 回答\n\n{answer}' if answer else ''
        docs = '\n'.join(f'### {r["title"]}\n{r["content"][:300]}...\n' for r in results[:3])

        content = f"""---
title: {query}
date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
type: qa
---

# {query}

## 搜索结果

{links}

{ai_section}

## 相关文档

{docs}
"""
        qa_file.write_text(content, encoding='utf-8')
        print(f"\n  ✅ 问答已沉淀到: {qa_file.relative_to(VAULT)}")


def search_wiki(query: str, top_k: int = 5) -> list[dict]:
    """搜索 wiki"""
    results = []
    query_lower = query.lower()

    for f in WIKI.rglob("*.md"):
        if f.name.startswith('.') or 'indexes' in str(f):
            continue

        try:
            content = f.read_text(encoding='utf-8', errors='ignore')
            content_lower = content.lower()

            if query_lower in content_lower:
                # 提取标题
                title_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
                title = title_match.group(1) if title_match else f.stem

                # 提取预览
                pos = content_lower.find(query_lower)
                preview = content[max(0, pos-30):pos+50]

                results.append({
                    "title": title,
                    "path": str(f.relative_to(VAULT)),
                    "content": content,
                    "preview": preview,
                    "score": content_lower.count(query_lower)
                })
        except:
            pass

    results.sort(key=lambda x: x["score"], reverse=True)
    return results[:top_k]


def extract_front_matter(content: str) -> dict:
    """提取 front matter 元数据"""
    import re
    fm = {}
    match = re.match(r'^---\n([\s\S]*?)\n---', content)
    if match:
        for line in match.group(1).split('\n'):
            if ':' in line:
                key, value = line.split(':', 1)
                fm[key.strip()] = value.strip()
    return fm


def cmd_compile(extract_concepts_flag: bool = True):
    """编译知识库 - Ingest 所有原始资料"""
    print_header("🔄 编译知识库 (Karpathy 方法)")

    if not RAW.exists():
        print("  ❌ raw/ 目录不存在")
        return

    # 1. 扫描原始文件
    raw_files = []
    for ext in SUPPORTED_EXTENSIONS:
        raw_files.extend(RAW.rglob(f"*{ext}"))
    print(f"  发现 {len(raw_files)} 个原始文件")

    # 2. 生成摘要
    summaries_dir = WIKI / "summaries"
    concepts_dir = WIKI / "concepts"
    summaries_dir.mkdir(parents=True, exist_ok=True)
    concepts_dir.mkdir(parents=True, exist_ok=True)

    compiled = 0
    concepts_found = set()

    for f in raw_files:
        if f.name.startswith('.'):
            continue

        summary_file = summaries_dir / f"{f.stem}.md"
        if not summary_file.exists():
            # 使用统一的解析函数读取内容
            content = parse_file(f)

            # 提取 front matter 元数据
            fm = extract_front_matter(content)

            # 优先使用 front matter 中的标题
            title = fm.get('title', f.stem)

            # 提取来源 URL（支持多种格式）
            source_url = fm.get('source', '') or fm.get('url', '') or fm.get('original-url', '')
            if not source_url and f.suffix == '.md':
                # 尝试从内容中查找 URL
                url_match = re.search(r'\[原文\]\((https?://[^\)]+)\)', content)
                if url_match:
                    source_url = url_match.group(1)

            # 提取概念
            concepts = extract_concepts(content)
            concepts_found.update(concepts)

            # tags
            tags_str = fm.get('tags', '')
            if not tags_str and concepts:
                tags_str = ', '.join(concepts[:3])

            concepts_list = '\n'.join(f'- {c}' for c in concepts[:5]) if concepts else '- 待提取'
            related_links = ', '.join(f'[[wiki/concepts/{c}]]' for c in concepts[:3]) if concepts else '无'

            # 生成摘要
            if source_url and source_url.startswith('http'):
                source_block = f"> 📎 原文：[{source_url}]({source_url})"
            else:
                source_block = f"> 原始文件：{f.name}"

            summary = f"""---
title: {title}
source: {source_url or str(f.relative_to(VAULT))}
date: {datetime.now().strftime('%Y-%m-%d')}
tags: [{tags_str}]
compiled: {datetime.now().isoformat()}
---

# {title}

{source_block}

## 摘要
（AI 补充中）

## 关键要点
{concepts_list}

## 相关概念
{related_links}
"""
            summary_file.write_text(summary, encoding='utf-8')
            compiled += 1

    print(f"  已生成 {compiled} 个摘要")

    # 3. 创建/更新概念页面
    if extract_concepts_flag and concepts_found:
        existing_concepts = [c.stem for c in concepts_dir.glob("*.md")]
        new_concepts = 0

        for concept in concepts_found:
            concept_file = concepts_dir / f"{concept}.md"
            if not concept_file.exists():
                # 找相关概念
                related = find_related_concepts(concept, existing_concepts)

                concept_content = f"""---
title: {concept}
type: concept
created: {datetime.now().strftime('%Y-%m-%d')}
related: [{', '.join(f'[[wiki/concepts/{r}]]' for r in related)}]
---

# {concept}

## 定义
（AI 补充中）

## 来源
{chr(10).join(f'- [[wiki/summaries/{f.stem}]]' for f in raw_files[:5])}

## 关联
{', '.join(f'[[wiki/concepts/{r}]]' for r in related) if related else '无'}
"""
                concept_file.write_text(concept_content, encoding='utf-8')
                new_concepts += 1

        print(f"  新增 {new_concepts} 个概念")

    # 4. 更新索引和日志
    update_index_file()
    update_log("compile", f"编译知识库，处理 {len(raw_files)} 个文件，新增 {compiled} 个摘要")

    print("  ✅ 编译完成")
    print("  运行 'python3 scripts/obsidian.py ai index' 查看索引")


def cmd_lint_enhanced():
    """增强版 Lint - 健康检查"""
    print_header("🔍 知识库健康检查")

    issues = []
    recommendations = []

    # 1. 检查孤立页面
    print("  检查孤立页面...")
    for f in WIKI.rglob("*.md"):
        if f.name.startswith('.') or 'indexes' in str(f):
            continue
        try:
            content = f.read_text(errors='ignore')
            links = re.findall(r'\[\[([^\]]+)\]\]', content)
            if not links:
                issues.append(f"孤立页面: {f.relative_to(VAULT)}")
        except:
            pass

    # 2. 检查未编译的 raw 文件
    print("  检查未编译文件...")
    raw_files = list(RAW.rglob("*.md"))
    summaries = list((WIKI / "summaries").rglob("*.md")) if (WIKI / "summaries").exists() else []
    if len(raw_files) > len(summaries):
        issues.append(f"未编译: {len(raw_files) - len(summaries)} 个原始文件")

    # 3. 检查概念完整性
    print("  检查概念网络...")
    concept_files = list((WIKI / "concepts").rglob("*.md")) if (WIKI / "concepts").exists() else []
    for cf in concept_files:
        try:
            content = cf.read_text(errors='ignore')
            if 'AI 补充中' in content or '待提取' in content:
                recommendations.append(f"概念待完善: {cf.stem}")
        except:
            pass

    # 4. 检查索引一致性
    print("  检查索引...")
    index_file = WIKI / "indexes" / "index.md"
    if index_file.exists():
        index_content = index_file.read_text()
        wiki_count = len(list(WIKI.rglob("*.md")))
        # 简单验证
        if "更新时间" not in index_content:
            issues.append("索引文件不完整")

    # 输出结果
    print()
    if issues:
        print("  ❌ 发现问题:")
        for issue in issues:
            print(f"     - {issue}")
    else:
        print("  ✅ 无问题")

    if recommendations:
        print("\n  💡 建议:")
        for rec in recommendations[:5]:
            print(f"     - {rec}")

    # 保存报告
    health_file = OUTPUTS / "health" / f"health-{datetime.now().strftime('%Y%m%d')}.md"
    health_file.write_text(f"""# 健康检查报告

**日期**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

## 问题 ({len(issues)} 个)

{chr(10).join(f'- {i}' for i in issues) if issues else '无'}

## 建议 ({len(recommendations)} 个)

{chr(10).join(f'- {r}' for r in recommendations) if recommendations else '无'}

## 统计

- 原始文件: {len(raw_files)}
- Wiki 页面: {len(list(WIKI.rglob('*.md'))) if WIKI.exists() else 0}
- 概念: {len(concept_files)}
- 摘要: {len(summaries)}
""")
    print(f"\n  📄 报告已保存: {health_file.relative_to(VAULT)}")


def cmd_update_index():
    """更新索引文件"""
    print_header("📋 更新索引")

    WIKI.mkdir(parents=True, exist_ok=True)
    (WIKI / "indexes").mkdir(parents=True, exist_ok=True)

    # 统计
    concepts = list((WIKI / "concepts").rglob("*.md")) if (WIKI / "concepts").exists() else []
    summaries = list((WIKI / "summaries").rglob("*.md")) if (WIKI / "summaries").exists() else []
    research = list((WIKI / "research").rglob("*.md")) if (WIKI / "research").exists() else []
    raw_files = list(RAW.rglob("*.md")) if RAW.exists() else []

    # 生成索引 - 避免 f-string 中使用 chr(10)
    concepts_section = '- 无概念'
    if concepts:
        items = '\n'.join(f'- [[wiki/concepts/{c.name}]]' for c in concepts[:10])
        concepts_section = items

    summaries_section = '- 无摘要'
    if summaries:
        items = '\n'.join(f'- [[wiki/summaries/{s.name}]] - {s.stem}' for s in summaries[:10])
        summaries_section = items

    raw_section = '- 无原始资料'
    if raw_files:
        items = '\n'.join(f'- [[{f.relative_to(VAULT)}]]' for f in raw_files[:10])
        raw_section = items

    content = f"""# Wiki Index

> East-lake 知识库索引 - 由 LLM 自动维护

## 统计

- 原始资料: {len(raw_files)} 篇
- 概念数: {len(concepts)} 个
- 摘要数: {len(summaries)} 篇
- 研究文档: {len(research)} 篇
- 更新时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

---

## 概念 (concepts/)

{concepts_section}

---

## 摘要 (summaries/)

{summaries_section}

---

## 原始资料 (raw/)

{raw_section}

---

*由 LLM 自动生成 - 最后更新：{datetime.now().strftime('%Y-%m-%d')}*
"""
    (WIKI / "indexes" / "index.md").write_text(content, encoding='utf-8')
    print("  ✅ 索引已更新")


# 别名 - 兼容旧代码
update_index_file = cmd_update_index


def update_log(action: str, detail: str):
    """更新日志"""
    LOG = WIKI / "indexes" / "log.md"
    entry = f"""
## [{datetime.now().strftime('%Y-%m-%d')}] {action}

- {detail}
"""
    if LOG.exists():
        existing = LOG.read_text()
        LOG.write_text(existing + entry)
    else:
        LOG.write_text(f"# Wiki Log\n{entry}")


# ============================================================
# 系统命令
# ============================================================

def cmd_system(args):
    """系统命令"""
    if args.action == "status":
        print_header("📊 系统状态")

        raw_count = len(list(RAW.rglob("*.md"))) if RAW.exists() else 0
        wiki_count = len(list(WIKI.rglob("*.md"))) if WIKI.exists() else 0
        concepts_count = len(list((WIKI / "concepts").rglob("*.md"))) if (WIKI / "concepts").exists() else 0
        summaries_count = len(list((WIKI / "summaries").rglob("*.md"))) if (WIKI / "summaries").exists() else 0

        print(f"  📂 原始资料 (raw/): {raw_count} 个")
        print(f"  📚 Wiki: {wiki_count} 个")
        print(f"     - 概念: {concepts_count}")
        print(f"     - 摘要: {summaries_count}")

        uncompiled = raw_count - summaries_count
        if uncompiled > 0:
            print(f"\n  ⚠️  未编译: {uncompiled} 个")
            print(f"     运行: python3 scripts/obsidian.py ai compile")

    elif args.action == "diagnose":
        run_py("diagnose.py")

    elif args.action == "install":
        print("📦 安装依赖...")
        packages = ["requests", "beautifulsoup4", "anthropic", "PyPDF2", "python-docx", "python-pptx"]
        for p in packages:
            subprocess.run([sys.executable, "-m", "pip", "install", p, "-q"], capture_output=True)
        print("  ✅ 完成")


# ============================================================
# 主程序
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description=f"East-lake v{VERSION} - Karpathy llm-wiki 知识库",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
核心操作 (Karpathy 方法):

  📥 Ingest (摄入):
    python obsidian.py ai compile      # 编译知识库（提取概念、建立链接）

  🔍 Query (查询):
    python obsidian.py ask <问题>      # 提问（答案自动沉淀到 wiki）
    python obsidian.py chat            # 交互对话

  🔧 Lint (维护):
    python obsidian.py ai lint         # 健康检查
    python obsidian.py ai index        # 更新索引

  📊 系统:
    python obsidian.py status          # 查看状态

示例:
  python obsidian.py ai compile        # 编译所有原始资料
  python obsidian.py ai lint           # 健康检查
  python obsidian.py ask "什么是RAG?"  # 提问
        """
    )

    parser.add_argument("command", nargs="?", help="主命令")
    parser.add_argument("action", nargs="?", help="子命令")
    parser.add_argument("args", nargs="?", help="参数")
    parser.add_argument("--tags", "-t", help="标签")
    parser.add_argument("--title", help="标题")
    parser.add_argument("--url", "-u", help="URL")

    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        print(f"\n版本: {VERSION}")
        return

    # 路由
    if args.command == "achieve" or args.command == "a":
        cmd_achieve(args)

    elif args.command == "web" or args.command == "w":
        cmd_web(args)

    elif args.command == "ask":
        # 支持 python obsidian.py ask "问题" 格式
        query = args.action or args.args or ""
        if query.startswith('"') or query.startswith("'"):
            query = query[1:-1]
        cmd_query(query)

    elif args.command == "chat":
        cmd_ai(argparse.Namespace(action="chat"))

    elif args.command == "ai":
        cmd_ai(argparse.Namespace(action=args.action or "stats"))

    elif args.command == "compile":
        cmd_compile()

    elif args.command == "lint":
        cmd_lint_enhanced()

    elif args.command == "index":
        cmd_update_index()

    elif args.command == "status":
        cmd_system(argparse.Namespace(action="status"))

    elif args.command == "install":
        cmd_system(argparse.Namespace(action="install"))

    else:
        print(f"未知命令: {args.command}")
        print("运行 python obsidian.py 查看帮助")


if __name__ == "__main__":
    main()